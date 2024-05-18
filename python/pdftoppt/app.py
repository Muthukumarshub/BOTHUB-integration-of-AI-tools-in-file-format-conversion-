from flask import Flask, render_template, request, send_file, jsonify
import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches

app = Flask(__name__)

def convert_pdf_to_ppt(pdf_file_path):
    try:
        ppt_file_path = pdf_file_path[:-4] + '.pptx'
        pdf_document = fitz.open(pdf_file_path)
        presentation = Presentation()

        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            image_matrix = fitz.Matrix(2, 2)  # Zoom factor
            pix = page.get_pixmap(matrix=image_matrix)
            image_path = f"page_{page_num}.png"
            pix.save(image_path)

            slide_layout = presentation.slide_layouts[5]  # Use a blank layout
            slide = presentation.slides.add_slide(slide_layout)
            left = top = Inches(1)
            pic = slide.shapes.add_picture(image_path, left, top)

            os.remove(image_path)  # Clean up the temporary image file

        presentation.save(ppt_file_path)
        return ppt_file_path, None
    except Exception as e:
        return None, str(e)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/convert', methods=['POST'])
def convert():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400

    pdf_file = request.files['file']
    if pdf_file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    if pdf_file.filename.endswith('.pdf'):
        pdf_file_path = os.path.join('uploads', pdf_file.filename)
        pdf_file.save(pdf_file_path)
        ppt_file_path, error = convert_pdf_to_ppt(pdf_file_path)
        if error:
            return jsonify({'error': error}), 500
        return jsonify({
            'download_link': f'/download/{os.path.basename(ppt_file_path)}'
        })
    else:
        return jsonify({'error': 'Please upload a PDF file'}), 400

@app.route('/download/<filename>')
def download(filename):
    ppt_file_path = os.path.join('uploads', filename)
    return send_file(ppt_file_path, as_attachment=True)

if __name__ == '__main__':
    if not os.path.exists('uploads'):
        os.makedirs('uploads')
    app.run(debug=True, port=5004)  # Set the port to 5004
